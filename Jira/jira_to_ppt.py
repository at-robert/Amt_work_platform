# coding=UTF-8
import sys
import os
import re
import time
import io
import platform
import pandas as pd

from pptx import Presentation
from pptx.util import Inches, Pt

from jira import JIRA
import json

from datetime import datetime

pwd = os.path.expanduser('~') + '/'


FILE_PATH_JIRA_CSV_MAC= pwd + 'Documents/Github/Amt_work_platform/Jira/CSV/'
FILE_PATH_JIRA_PIC_MAC= pwd + 'Documents/Github/Amt_work_platform/Jira/PIC/'
FILE_PATH_JIRA_PPT_MAC= pwd + 'Documents/Github/Amt_work_platform/Jira/PPT/'

FILE_PATH_JIRA_CSV=r"D:\work_platform\Github\Amt_work_platform\Jira\CSV\\"
FILE_PATH_JIRA_PIC=r"D:\work_platform\Github\Amt_work_platform\Jira\PIC\\"
FILE_PATH_JIRA_PPT=r"D:\work_platform\Github\Amt_work_platform\Jira\PPT\\"

#----------------------------------------------------------------------
def change_table_font_size(prs_):
    for slide in prs_.slides:
    # 遍历每个形状
        for shape in slide.shapes:
            # 检查形状是否为表格
            if shape.has_table:
                # 遍历每个单元格
                for row in shape.table.rows:
                    for cell in row.cells:
                        # 修改单元格中文本框的字体大小
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                font = run.font
                                font.size = Pt(12)

#----------------------------------------------------------------------
def add_table(shapes_,number_):

    item_ = ['OSD 操作介面','SYS 系統行為', 'Audio 爆音異音', 'PQ 畫質相關', 'Video 畫異、閃屏','待處理', 'Total' ]
    cols = len(item_)
    rows = 2
    left = Inches(1.0)
    top = Inches(2.0)
    width = Inches(3.0)
    height = Inches(0.4)

    table = shapes_.add_table(rows, cols, left, top, width, height).table

    # set column widths
    # table.columns[0].width = Inches(2.0)
    # table.columns[1].width = Inches(4.0)

    for i in range(0,len(item_)):
        if(item_[i] == 'OSD 操作介面'):
            table.columns[i].width = Inches(0.9)
        elif(item_[i] == 'SYS 系統行為'):
            table.columns[i].width = Inches(0.9)
        elif(item_[i] == 'Audio 爆音異音'):
            table.columns[i].width = Inches(0.9)
        elif(item_[i] == 'PQ 畫質相關'):
            table.columns[i].width = Inches(0.9)
        elif(item_[i] == '待處理'):
            table.columns[i].width = Inches(0.8)
        elif(item_[i] == 'Total'):
            table.columns[i].width = Inches(0.6)
        else:
            table.columns[i].width = Inches(1.0)

        table.cell(0, i).text = item_[i]

    for i in range(0,len(item_)):
        table.cell(1, i).text = str(number_[i])


#----------------------------------------------------------------------
def add_slide(prs, layout, title, number_,img_name_):
    """Return slide newly added to `prs` using `layout` and having `title`."""
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    shape = slide.shapes

    add_table(shape,number_)

    left = Inches(3)
    top = Inches(4)
    height = Inches(5)
    
    if (platform.system() == 'Darwin'):
        slide.shapes.add_picture(FILE_PATH_JIRA_PIC_MAC + img_name_, left, top, height)
    else:
        slide.shapes.add_picture(FILE_PATH_JIRA_PIC + img_name_, left, top, height)

    return slide


#----------------------------------------------------------------------
def check_if_png_exist(prj_name_):
    if (platform.system() == 'Darwin'):
        return os.path.exists(FILE_PATH_JIRA_PIC_MAC + prj_name_ + '.png')
    else:
        return os.path.exists(FILE_PATH_JIRA_PIC + prj_name_ + '.png')


#----------------------------------------------------------------------
if __name__ == "__main__":

    if (platform.system() == 'Darwin'):
        df = pd.read_csv(FILE_PATH_JIRA_CSV_MAC + 'Jira_Statistic_.csv')
    else:
        df = pd.read_csv(FILE_PATH_JIRA_CSV + 'Jira_Statistic_.csv')

    # print(df)

    prs = Presentation()
    title_only_slide_layout = prs.slide_layouts[5]


    for index, row in df.iterrows():
        prj_name = row['Model']
        list_ = row.values.flatten().tolist()[1:]

        # Only .png exist would generate the report
        if ( check_if_png_exist(prj_name) ):
            add_slide(prs, title_only_slide_layout, prj_name + time.strftime("%c"), list_, prj_name + '.png')

    change_table_font_size(prs)


    if (platform.system() == 'Darwin'):
        prs.save(FILE_PATH_JIRA_PPT_MAC + 'jira_npi.pptx')
    else:
        prs.save(FILE_PATH_JIRA_PPT + 'jira_npi.pptx')